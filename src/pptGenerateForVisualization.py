from pptx import Presentation #
import collections
import collections.abc
from pptx.util import Inches
import os
from PIL import Image
import glob
import numpy as np

def generatePPT(dirName:str, targetFileName:str):
    try:
        from collections.abc import Container
    except ImportError:
        from collections import Container
    collections.Container = collections.abc.Container
    collections.Mapping = collections.abc.Mapping
    collections.MutableMapping = collections.abc.MutableMapping
    collections.Iterable = collections.abc.Iterable
    collections.MutableSet = collections.abc.MutableSet
    collections.Callable = collections.abc.Callable

    # current_dir = os.getcwd() + "/"
    imgarr = [] #later to be array of images
    image_dir = dirName+"/"


    #Makes the array of images into imgarr
    for filename in os.listdir(image_dir):
        if filename.endswith(".jpg") or filename.endswith(".png"):
            imgarr.append(os.path.join(image_dir, filename))
    imgarr.sort()
    print(imgarr)
    #Makes the slides in powerpoint
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for numslides in range(len(imgarr)):
        slide = prs.slides.add_slide(blank_slide_layout)
        left = Inches(-.5)
        # top = Inches(-1.2)
        top = Inches(0)
        pic = slide.shapes.add_picture(imgarr[numslides], left, top, Inches(11))


    prs.save(targetFileName+".pptx")